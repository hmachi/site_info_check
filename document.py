import shutil
import datetime
from docx import Document
from pptx.util import Inches


def create_document(site_info, whois, ip, geolocation, url, domain, nowStr):

    new_doc_name = domain + "_" + nowStr + ".docx"
    new_doc_path = "document/" + new_doc_name

    shutil.copyfile("document/document.docx", new_doc_path)

    doc = Document(new_doc_path)

    now = datetime.datetime.now(datetime.timezone(
        datetime.timedelta(hours=9)))

    # 行単位の置換
    for paragraph in doc.paragraphs:
        replace_text(paragraph, whois, ip, geolocation, now, url, domain)

        # 入力したURLの画像を設定する
        if paragraph.text == '${site_screenshot}':
            paragraph.text = ""
            r = paragraph.add_run()
            r.add_picture(site_info['screenshot_path'], width=Inches(6.2))

    paragraphs = (paragraph
                  for table in doc.tables
                  for row in table.rows
                  for cell in row.cells
                  for paragraph in cell.paragraphs)

    # テーブルのセル単位の置換
    for paragraph in paragraphs:
        replace_text(paragraph, whois, ip, geolocation, now, url, domain)

    # 画像の貼り付け
    for index, screenshot_path in enumerate(ip['screenshot_path_list']):
        doc.add_picture(screenshot_path, width=Inches(6.2))
        if index % 2 == 0:
            doc.add_paragraph()

    # 改ページ
    doc.add_page_break()

    for index, screenshot_path in enumerate(geolocation['screenshot_path_list']):
        doc.add_picture(screenshot_path, width=Inches(6.2))
        if index % 2 == 0:
            doc.add_paragraph()

    doc.save(new_doc_path)

    return {
        "name": new_doc_name,
        "path": new_doc_path
    }


def replace_text(paragraph, whois, ip, geolocation, now, url, domain):
    """
    word文書の文字の置換
    ※行単位・セル単位で文字の置換をする必要がある
    """

    try:
        if paragraph.text == "${doc_creation_date}":
            paragraph.text = "令和" + \
                str(int(str(now.strftime('%Y'))[1:]) - 18) + "年" + \
                now.strftime('%m') + "月" + now.strftime('%d') + "日"

        if paragraph.text == "${name}":
            paragraph.text = ip["company"] if ip["company"] != "" else "不明"
        if paragraph.text == "${url}":
            paragraph.text = url
        if paragraph.text == "${url_brackets}":
            paragraph.text = "【" + url + "】"
        if paragraph.text == "${address}":
            paragraph.text = ip["address"] if ip["address"] != "" else "不明"
        if paragraph.text == "${domain}":
            paragraph.text = domain
        if paragraph.text == "${host}":
            paragraph.text = "ホスト名：" + \
                (ip["host"] if ip["host"] != "" else "不明")
        if paragraph.text == "${area}":
            paragraph.text = "エリア：" + \
                (ip["area"]
                 if ip["area"] != "" else "不明")
        if paragraph.text == "${ipv4}":
            paragraph.text = ip["ipv4"] if ip["ipv4"] != "" else "不明"

        if paragraph.text == "${domain_name}":
            paragraph.text = whois["domain_name"] if whois["domain_name"] != "" else "不明"
        if paragraph.text == "${expiration_date}":
            paragraph.text = whois["expiration_date"] if whois["expiration_date"] != "" else "不明"
        if paragraph.text == "${domain_status}":
            paragraph.text = whois["domain_status"] if whois["domain_status"] != "" else "不明"
        if paragraph.text == "${name_server}":
            paragraph.text = whois["name_server"] if whois["name_server"] != "" else "不明"
        if paragraph.text == "${registrant_email}":
            paragraph.text = whois["registrant_email"] if whois["registrant_email"] != "" else "不明"
        if paragraph.text == "${admin_email}":
            paragraph.text = whois["admin_email"] if whois["admin_email"] != "" else "不明"
        if paragraph.text == "${creation_date}":
            paragraph.text = whois["creation_date"] if whois["creation_date"] != "" else "不明"
        if paragraph.text == "${registrar}":
            paragraph.text = whois["registrar"] if whois["registrar"] != "" else "不明"

    except Exception as e:
        print(e)
